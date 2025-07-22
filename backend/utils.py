import hashlib
import os
import pdfplumber
import docx
from pptx import Presentation
from typing import List

# No path changes needed in utils.py itself, as its functions operate on paths passed to them.
# The calling modules are responsible for constructing paths relative to 'persistent_storage'.

def compute_md5(file_path: str) -> str:
    """Compute the MD5 hash of a file to detect duplicates."""
    hash_md5 = hashlib.md5()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()

def compute_text_md5(text: str) -> str:
    """Compute MD5 hash of a string (e.g., for scraped web content)."""
    return hashlib.md5(text.encode('utf-8')).hexdigest()

def compute_md5_from_text(text: str) -> str:
    """Compute the MD5 hash of a text string to detect duplicates."""
    return hashlib.md5(text.encode("utf-8")).hexdigest()

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file using pdfplumber."""
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
    except Exception as e:
        print(f"[PDF Error] {file_path}: {e}")
    return text

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a DOCX file using python-docx."""
    text = ""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"[DOCX Error] {file_path}: {e}")
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX file using python-pptx."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"[PPTX Error] {file_path}: {e}")
    return text

def extract_text_from_file(file_path: str) -> str:
    """Determine the file type and extract text accordingly."""
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        print(f"[Unsupported Format] {file_path}")
        return ""

def ensure_directory(path: str):
    """Create a directory if it doesn't exist."""
    os.makedirs(path, exist_ok=True)

def split_text_into_chunks(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """
    Splits long text into overlapping chunks for vector embedding.
    Example: chunk_size=1000, overlap=200 means chunk1=0-1000, chunk2=800-1800, etc.
    """
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end]
        chunks.append(chunk)
        start += chunk_size - overlap
    return chunks